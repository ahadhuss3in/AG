"""This file is for microsoft office files, the ones ending in .docx or .pptx.

 use python-docx for word documents and python-pptx for powerpoint
Which one runs depends on the file's ending. Word documents can be para 
or table . Powerpoint files are read slide by
slide,
"""

from pathlib import Path

import logfire
from docx import Document as DocxDocument
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from . import LoadedDocument, check_file_size, warn_if_empty


def _format_table(rows: list[list[str]]) -> str:
    """Turn a table's rows into plain text, one row per line."""
    return "\n".join(" | ".join(cell for cell in row) for row in rows)


def _load_docx(file_path: Path) -> tuple[str, dict]:
    doc = DocxDocument(file_path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    text_parts = list(paragraphs)

    #for table go through seperate
    for table_index, table in enumerate(doc.tables, start=1):
        rows = [[cell.text for cell in row.cells] for row in table.rows]
        table_text = _format_table(rows)
        if table_text.strip():
            text_parts.append(f"[table {table_index}]\n{table_text}")

    text = "\n".join(text_parts)
    metadata = {
        "file_type": "docx",
        "paragraph_count": len(paragraphs),
        "table_count": len(doc.tables),
    }
    return text, metadata


def _iter_shapes(shapes):
    """Go through every shape on a slide.
    """
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def _load_pptx(file_path: Path) -> tuple[str, dict]:
    presentation = Presentation(file_path)
    slides_text = []
    notes_found = 0
    tables_found = 0

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_parts = []

        for shape in _iter_shapes(slide.shapes):
            if shape.has_text_frame and shape.text_frame.text.strip():
                slide_parts.append(shape.text_frame.text)

            if shape.has_table:
                tables_found += 1
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                table_text = _format_table(rows)
                if table_text.strip():
                    slide_parts.append(table_text)

        # check for speaker notes.
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text
            if notes_text.strip():
                notes_found += 1
                slide_parts.append(f"speaker notes: {notes_text}")

        if slide_parts:
            slides_text.append(f"[slide {slide_number}]\n" + "\n".join(slide_parts))

    text = "\n\n".join(slides_text)
    metadata = {
        "file_type": "pptx",
        "slide_count": len(presentation.slides),
        "notes_found": notes_found,
        "tables_found": tables_found,
    }
    return text, metadata


_HANDLERS = {
    ".docx": _load_docx,
    ".pptx": _load_pptx,
}


def loadoffice(file_path: str | Path) -> LoadedDocument:
    """Read a .docx or .pptx file into a LoadedDocument.

    Raises a clear error for any other file ending, since this loader only
    knows how to handle these two office formats.
    """
    file_path = Path(file_path)
    suffix = file_path.suffix.lower()

    handler = _HANDLERS.get(suffix)
    if handler is None:
        raise ValueError(f"office_loader can't handle extension: {suffix}")

    with logfire.span("load_office", file=str(file_path), file_type=suffix):
        try:
            check_file_size(file_path)
            text, metadata = handler(file_path)
        except Exception:
            logfire.exception("could not read this office file", file=str(file_path))
            raise

        warn_if_empty(text, file_path)

        logfire.info(
            "read office file",
            file=str(file_path),
            chars=len(text),
            **metadata,
        )

        return LoadedDocument(text=text, source=str(file_path), metadata=metadata)
